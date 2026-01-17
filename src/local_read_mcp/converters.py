################################################################################
# Local Read MCP - Document Converters
# Extracted from MiroThinker, excluding OpenAI dependencies
################################################################################

import json
import logging
import os
import re
import shutil
import tempfile
import time
import traceback
import csv as csv_module
from typing import Any, Union, Optional, Dict, List
from urllib.parse import quote, unquote, urlparse, urlunparse

# Optional imports
try:
    import mammoth
except ImportError:
    mammoth = None

try:
    import markdownify
except ImportError:
    markdownify = None

try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    openpyxl = None
    get_column_letter = None

try:
    import pdfminer
    import pdfminer.high_level
except ImportError:
    pdfminer = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    from markitdown import MarkItDown
except ImportError:
    MarkItDown = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import yaml
except ImportError:
    yaml = None

try:
    import html
except ImportError:
    html = None

# Media file extension constants
IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "gif", "webp"}
AUDIO_EXTENSIONS = {"wav", "mp3", "m4a"}
VIDEO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "webm"}
MEDIA_EXTENSIONS = IMAGE_EXTENSIONS | AUDIO_EXTENSIONS | VIDEO_EXTENSIONS

class DocumentConverterResult:
    """Document conversion result with enhanced metadata and structure.

    This class encapsulates the result of converting a document to text/markdown format,
    including optional metadata, sections, tables, images, and pagination information.

    Attributes:
        title: Document title (optional, may be None)
        text_content: Converted text/markdown content
        metadata: Additional metadata dict (file size, timestamps, etc.)
        sections: List of extracted sections with headings
        tables: List of extracted table information
        images: List of extracted image information (for PDFs with extract_images=True)
        pagination_info: Pagination details (page count, offsets, etc.)
        processing_time_ms: Processing time in milliseconds (optional)
        error: Error message if conversion failed (optional)

    Example:
        >>> result = DocumentConverterResult(
        ...     title="My Document",
        ...     text_content="# Heading\\n\\nContent here",
        ...     metadata={"file_size": 12345}
        ... )
        >>> print(result.title)
        My Document
    """

    def __init__(
        self,
        title: Union[str, None] = None,
        text_content: str = "",
        metadata: Optional[Dict[str, Any]] = None,
        sections: Optional[List[Dict[str, Any]]] = None,
        tables: Optional[List[Dict[str, Any]]] = None,
        images: Optional[List[Dict[str, Any]]] = None,
        pagination_info: Optional[Dict[str, Any]] = None,
        processing_time_ms: Optional[int] = None,
        error: Optional[str] = None
    ):
        """Initialize document converter result.

        Args:
            title: Document title (optional)
            text_content: Converted text content
            metadata: Additional metadata (default: empty dict)
            sections: List of document sections (default: empty list)
            tables: List of extracted tables (default: empty list)
            images: List of extracted images (default: empty list)
            pagination_info: Pagination information (default: empty dict)
            processing_time_ms: Processing time in milliseconds (optional)
            error: Error message if failed (optional)
        """
        self.title: Union[str, None] = title
        self.text_content: str = text_content
        self.metadata = metadata or {}
        self.sections = sections or []
        self.tables = tables or []
        self.images = images or []
        self.pagination_info = pagination_info or {}
        self.processing_time_ms = processing_time_ms
        self.error = error

    def to_dict(self) -> Dict[str, Any]:
        """Convert result to dictionary for JSON serialization.

        Returns:
            Dictionary representation of the conversion result.

        Example:
            >>> result = DocumentConverterResult(title="Test", text_content="Content")
            >>> d = result.to_dict()
            >>> d["title"]
            'Test'
        """
        result = {
            "title": self.title,
            "text_content": self.text_content,
            "metadata": self.metadata,
            "sections": self.sections,
            "tables": self.tables,
            "pagination_info": self.pagination_info,
        }
        if self.processing_time_ms is not None:
            result["processing_time_ms"] = self.processing_time_ms
        if self.error is not None:
            result["error"] = self.error
        return result

class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in [
                    "http",
                    "https",
                    "file",
                ]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(
                    parsed_url._replace(path=quote(unquote(parsed_url.path)))
                )  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


# Removed duplicate DocumentConverterResult class definition (was at line 198-204)


def convert_html_to_md(html_content: str) -> DocumentConverterResult:
    """Convert HTML content to Markdown format.

    This function parses HTML content using BeautifulSoup, removes script and style tags,
    and converts the remaining content to Markdown format using a custom Markdown converter.

    Args:
        html_content: Raw HTML content string to convert

    Returns:
        DocumentConverterResult with the converted Markdown text and title

    Raises:
        Exception: If HTML parsing or conversion fails

    Example:
        >>> html = "<html><head><title>Test</title></head><body><h1>Hello</h1></body></html>"
        >>> result = convert_html_to_md(html)
        >>> "Hello" in result.text_content
        True

    Note:
        - JavaScript and CSS are automatically removed
        - Only the <body> content is converted if present, otherwise full document
        - Uses _CustomMarkdownify for conversion with enhanced security
    """
    soup = BeautifulSoup(html_content, "html.parser")
    for script in soup(["script", "style"]):
        script.extract()

    # Print only the main content
    body_elm = soup.find("body")
    webpage_text = ""
    if body_elm:
        webpage_text = _CustomMarkdownify().convert_soup(body_elm)
    else:
        webpage_text = _CustomMarkdownify().convert_soup(soup)

    assert isinstance(webpage_text, str)

    return DocumentConverterResult(
        title=None if soup.title is None else soup.title.string,
        text_content=webpage_text,
    )


def HtmlConverter(
    local_path: str,
    extract_metadata: bool = False,
    extract_sections: bool = False
) -> DocumentConverterResult:
    """
    Convert an HTML file to Markdown format with enhanced features.

    Args:
        local_path: Path to the HTML file to convert.
        extract_metadata: Whether to extract metadata (file size, etc.)
        extract_sections: Whether to extract sections from content

    Returns:
        DocumentConverterResult containing the converted Markdown text and optional metadata/sections.
    """
    try:
        with open(local_path, "rt", encoding="utf-8") as fh:
            html_content = fh.read()

        # Convert HTML to markdown
        soup = BeautifulSoup(html_content, "html.parser")
        for script in soup(["script", "style"]):
            script.extract()

        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        # Apply content limit
        webpage_text = apply_content_limit(webpage_text)

        # Prepare metadata
        metadata = {}
        sections = []

        if extract_metadata:
            metadata = {
                "file_path": local_path,
                "file_size": os.path.getsize(local_path) if os.path.exists(local_path) else None,
                "file_extension": os.path.splitext(local_path)[1],
                "conversion_timestamp": time.time()
            }

        if extract_sections:
            sections = extract_sections_from_markdown(webpage_text)

        # Get title from HTML or use filename
        title = None
        if soup.title and soup.title.string:
            title = soup.title.string
        else:
            # Use filename without extension as fallback
            filename = os.path.basename(local_path)
            title = os.path.splitext(filename)[0]

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
            metadata=metadata,
            sections=sections,
            tables=[],  # HTML tables not extracted in basic version
            processing_time_ms=None
        )

    except Exception as e:
        return DocumentConverterResult(
            title=None,
            text_content=f"Error converting HTML: {str(e)}",
            error=str(e)
        )


def DocxConverter(
    local_path: str,
    extract_metadata: bool = False,
    extract_sections: bool = False
) -> DocumentConverterResult:
    """
    Convert a DOCX file to Markdown format with enhanced features.

    Uses mammoth library to first convert DOCX to HTML, then converts
    the HTML to Markdown.

    Args:
        local_path: Path to the DOCX file to convert.
        extract_metadata: Whether to extract metadata (file size, etc.)
        extract_sections: Whether to extract sections from content

    Returns:
        DocumentConverterResult containing the converted Markdown text and optional metadata/sections.
    """
    if mammoth is None:
        return DocumentConverterResult(
            title=None,
            text_content="[Error: mammoth not installed]",
            error="mammoth not installed"
        )

    try:
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value

        # Convert HTML to markdown
        soup = BeautifulSoup(html_content, "html.parser")
        for script in soup(["script", "style"]):
            script.extract()

        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        # Apply content limit
        webpage_text = apply_content_limit(webpage_text)

        # Prepare metadata
        metadata = {}
        sections = []

        if extract_metadata:
            metadata = {
                "file_path": local_path,
                "file_size": os.path.getsize(local_path) if os.path.exists(local_path) else None,
                "file_extension": os.path.splitext(local_path)[1],
                "conversion_timestamp": time.time()
            }

        if extract_sections:
            sections = extract_sections_from_markdown(webpage_text)

        # Get title from HTML or use filename
        title = None
        if soup.title and soup.title.string:
            title = soup.title.string
        else:
            # Use filename without extension as fallback
            filename = os.path.basename(local_path)
            title = os.path.splitext(filename)[0]

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
            metadata=metadata,
            sections=sections,
            tables=[],  # DOCX tables not extracted in basic version
            processing_time_ms=None
        )

    except Exception as e:
        return DocumentConverterResult(
            title=None,
            text_content=f"Error converting DOCX: {str(e)}",
            error=str(e)
        )


def XlsxConverter(
    local_path: str,
    extract_metadata: bool = False,
    extract_tables: bool = False
) -> DocumentConverterResult:
    """
    Converts Excel files to Markdown using openpyxl with enhanced features.
    Preserves color formatting and other cell styling information.

    Args:
        local_path: Path to the Excel file
        extract_metadata: Whether to extract metadata (file size, sheet info, etc.)
        extract_tables: Whether to extract table information (currently always extracts tables)

    Returns:
        DocumentConverterResult with the Markdown representation and optional metadata
    """
    # Load the workbook
    wb = openpyxl.load_workbook(local_path, data_only=True)
    md_content = ""

    # Prepare metadata
    metadata = {}
    tables = []

    if extract_metadata:
        # Collect workbook metadata
        sheet_names = wb.sheetnames
        metadata = {
            "file_path": local_path,
            "file_size": os.path.getsize(local_path) if os.path.exists(local_path) else None,
            "file_extension": os.path.splitext(local_path)[1],
            "conversion_timestamp": time.time(),
            "sheet_count": len(sheet_names),
            "sheet_names": sheet_names,
            "active_sheet": wb.active.title if wb.active else None
        }

    if extract_tables:
        # For Excel, each sheet is considered a table
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Get table dimensions
            min_row, min_col = 1, 1
            max_row = max(
                (cell.row for cell in sheet._cells.values() if cell.value is not None),
                default=0,
            )
            max_col = max(
                (cell.column for cell in sheet._cells.values() if cell.value is not None),
                default=0,
            )
            if max_row > 0 and max_col > 0:
                tables.append({
                    "sheet_name": sheet_name,
                    "rows": max_row,
                    "columns": max_col,
                    "has_data": True
                })

    # Helper function to convert RGB color to hex
    def rgb_to_hex(rgb_value):
        if not rgb_value:
            return None

        # Convert RGB value to string for processing
        rgb_string = str(rgb_value)

        # Handle RGB format like 'RGB(255, 255, 255)'
        if isinstance(rgb_value, str) and rgb_string.startswith("RGB"):
            rgb_match = re.match(r"RGB\((\d+), (\d+), (\d+)\)", rgb_string)
            if rgb_match:
                r, g, b = map(int, rgb_match.groups())
                return f"#{r:02x}{g:02x}{b:02x}"

        # Special handling for FFFFFFFF (white) and 00000000 (transparent/none)
        if rgb_string in ["FFFFFFFF", "00000000", "none", "auto"]:
            return None

        # Handle ARGB format (common in openpyxl)
        if len(rgb_string) == 8:  # ARGB format like 'FF5733FF'
            return f"#{rgb_string[2:]}"  # Strip alpha channel

        # Handle direct hex values like 'FF5733'
        if isinstance(rgb_value, str):
            return f"#{rgb_string}" if not rgb_string.startswith("#") else rgb_string

        return None  # Return None for unrecognized formats

    # Helper function to detect and format cell styling
    def get_cell_format_info(cell):
        info = {}

        # Get background color if it exists
        if cell.fill and hasattr(cell.fill, "fgColor") and cell.fill.fgColor:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.fill.fgColor, "rgb") and cell.fill.fgColor.rgb:
                rgb_value = cell.fill.fgColor.rgb
            elif hasattr(cell.fill.fgColor, "value") and cell.fill.fgColor.value:
                rgb_value = cell.fill.fgColor.value

            if rgb_value:
                bg_color = rgb_to_hex(rgb_value)
                if bg_color:  # Skip transparent or white (handled in rgb_to_hex)
                    info["bg_color"] = bg_color

        # Get font color if it exists
        if cell.font and hasattr(cell.font, "color") and cell.font.color:
            # Get the RGB value - in openpyxl this can be stored in different attributes
            rgb_value = None
            if hasattr(cell.font.color, "rgb") and cell.font.color.rgb:
                rgb_value = cell.font.color.rgb
            elif hasattr(cell.font.color, "value") and cell.font.color.value:
                rgb_value = cell.font.color.value

            if rgb_value:
                font_color = rgb_to_hex(rgb_value)
                if font_color:  # Skip transparent (handled in rgb_to_hex)
                    info["font_color"] = font_color

        # Get font weight (bold)
        if cell.font and cell.font.bold:
            info["bold"] = True

        # Get font style (italic)
        if cell.font and cell.font.italic:
            info["italic"] = True

        # Get font underline
        if cell.font and cell.font.underline and cell.font.underline != "none":
            info["underline"] = True

        return info

    # Process each sheet in the workbook
    for sheet_name in wb.sheetnames:
        try:
            sheet = wb[sheet_name]
            md_content += f"## {sheet_name}\n\n"

            # Get the dimensions of the used part of the sheet
            min_row, min_col = 1, 1
            max_row = max(
                (cell.row for cell in sheet._cells.values() if cell.value is not None),
                default=0,
            )
            max_col = max(
                (
                    cell.column
                    for cell in sheet._cells.values()
                    if cell.value is not None
                ),
                default=0,
            )

            if max_row == 0 or max_col == 0:
                md_content += "This sheet is empty.\n\n"
                continue
        except Exception as e:
            error_msg = f"Error processing sheet '{sheet_name}': {str(e)}"
            print(error_msg)
            md_content += (
                f"## {sheet_name}\n\nError processing this sheet: {str(e)}\n\n"
            )
            continue

        try:
            # First, determine column widths
            col_widths = {}
            for col_idx in range(min_col, max_col + 1):
                max_length = 0
                for row_idx in range(min_row, max_row + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""
                        max_length = max(max_length, len(cell_value))
                    except Exception as e:
                        print(
                            f"Warning: Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        max_length = max(max_length, 10)  # Use reasonable default
                col_widths[col_idx] = max(max_length + 2, 5)  # Min width of 5

            # Start building the table
            # Header row with column separators
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += " " + " " * col_widths[col_idx] + " |"
            md_content += "\n"

            # Separator row
            md_content += "|"
            for col_idx in range(min_col, max_col + 1):
                md_content += ":" + "-" * col_widths[col_idx] + ":|"
            md_content += "\n"

            # Data rows
            for row_idx in range(min_row, max_row + 1):
                md_content += "|"
                for col_idx in range(min_col, max_col + 1):
                    try:
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        cell_value = str(cell.value) if cell.value is not None else ""

                        # Get formatting info
                        try:
                            format_info = get_cell_format_info(cell)
                        except Exception as e:
                            print(
                                f"Warning: Error getting formatting for cell at row {row_idx}, column {col_idx}: {str(e)}"
                            )
                            format_info = {}

                        formatted_value = cell_value

                        # Add HTML-style formatting if needed
                        if format_info:
                            style_parts = []

                            if "bg_color" in format_info:
                                style_parts.append(
                                    f"background-color:{format_info['bg_color']}"
                                )

                            if "font_color" in format_info:
                                style_parts.append(f"color:{format_info['font_color']}")

                            span_attributes = []
                            if style_parts:
                                span_attributes.append(
                                    f'style="{"; ".join(style_parts)}"'
                                )

                            # Format with bold/italic/underline if needed
                            inner_value = cell_value
                            if "bold" in format_info:
                                inner_value = f"<strong>{inner_value}</strong>"
                            if "italic" in format_info:
                                inner_value = f"<em>{inner_value}</em>"
                            if "underline" in format_info:
                                inner_value = f"<u>{inner_value}</u>"

                            # Only add a span if we have style attributes
                            if span_attributes:
                                formatted_value = f"<span {' '.join(span_attributes)}>{inner_value}</span>"
                            else:
                                formatted_value = inner_value

                        # Pad to column width and add to markdown
                        padding = col_widths[col_idx] - len(cell_value)
                        padded_value = " " + formatted_value + " " * (padding + 1)
                        md_content += padded_value + "|"
                    except Exception as e:
                        print(
                            f"Error processing cell at row {row_idx}, column {col_idx}: {str(e)}"
                        )
                        # Add a placeholder for the failed cell
                        padded_value = " [Error] " + " " * (col_widths[col_idx] - 7)
                        md_content += padded_value + " |"

                md_content += "\n"
        except Exception as e:
            error_msg = f"Error generating table for sheet '{sheet_name}': {str(e)}\n{traceback.format_exc()}"
            print(error_msg)
            md_content += f"Error generating table: {str(e)}\n\n"

        # Add formatting legend
        has_formatting = False
        for row_idx in range(min_row, max_row + 1):
            for col_idx in range(min_col, max_col + 1):
                cell = sheet.cell(row=row_idx, column=col_idx)
                if get_cell_format_info(cell):
                    has_formatting = True
                    break
            if has_formatting:
                break

        if has_formatting:
            md_content += "\n### Formatting Information\n"
            md_content += "The table above includes HTML formatting to represent colors and styles from the original Excel file.\n"
            md_content += "This formatting may not display in all Markdown viewers.\n"

        md_content += "\n\n"  # Extra newlines between sheets

    # Apply content limit
    final_content = apply_content_limit(md_content.strip())

    # Use filename without extension as title
    filename = os.path.basename(local_path)
    title = os.path.splitext(filename)[0]

    return DocumentConverterResult(
        title=title,
        text_content=final_content,
        metadata=metadata,
        sections=[],  # Excel doesn't have sections in markdown sense
        tables=tables,
        processing_time_ms=None  # Can be calculated by caller
    )


def PptxConverter(local_path: str) -> DocumentConverterResult:
    """
    Converts PPTX files to Markdown. Supports headings, tables and images with alt text.

    Args:
        local_path: Path to the PPTX file

    Returns:
        DocumentConverterResult containing the converted Markdown text
    """

    def is_picture(shape):
        """Check if a shape is a picture"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def is_table(shape):
        """Check if a shape is a table"""
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    if not local_path.endswith(".pptx"):
        return DocumentConverterResult(
            title=None,
            text_content=f"Error: Expected .pptx file, got: {local_path}",
        )

    md_content = ""
    presentation = pptx.Presentation(local_path)
    slide_num = 0

    for slide in presentation.slides:
        slide_num += 1
        md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"
        title = slide.shapes.title

        for shape in slide.shapes:
            # Pictures
            if is_picture(shape):
                # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                alt_text = ""
                try:
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                except Exception:
                    pass
                # A placeholder name
                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += (
                    "\n!["
                    + (alt_text if alt_text else shape.name)
                    + "]("
                    + filename
                    + ")\n"
                )

            # Tables
            if is_table(shape):
                html_table = "<html><body><table>"
                first_row = True
                for row in shape.table.rows:
                    html_table += "<tr>"
                    for cell in row.cells:
                        if first_row:
                            html_table += "<th>" + html.escape(cell.text) + "</th>"
                        else:
                            html_table += "<td>" + html.escape(cell.text) + "</td>"
                    html_table += "</tr>"
                    first_row = False
                html_table += "</table></body></html>"

                # Note: This would require a separate HTML to Markdown converter function
                # In this version, I'm assuming a convert_html_to_md function exists
                md_content += (
                    "\n" + convert_html_to_md(html_table).text_content.strip() + "\n"
                )

            # Text areas
            elif shape.has_text_frame:
                if shape == title:
                    md_content += "# " + shape.text.lstrip() + "\n"
                else:
                    md_content += shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                md_content += notes_frame.text
            md_content = md_content.strip()

    return DocumentConverterResult(
        title=None,
        text_content=md_content.strip(),
    )


def ZipConverter(local_path: str, **kwargs):
    """
    Extracts ZIP files to a temporary directory and processes each file according to its extension.
    Returns a combined result of all processed files.
    """
    import zipfile

    temp_dir = tempfile.mkdtemp(prefix="zip_extract_")
    md_content = f"# Extracted from ZIP: {os.path.basename(local_path)}\n\n"

    try:
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            zip_ref.extractall(temp_dir)

        # Get all extracted files
        extracted_files = []
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                rel_path = os.path.relpath(file_path, temp_dir)
                extracted_files.append((file_path, rel_path))

        if not extracted_files:
            md_content += "The ZIP file is empty or contains no files.\n"
        else:
            md_content += f"Total files extracted: {len(extracted_files)}\n\n"

            for file_path, rel_path in extracted_files:
                md_content += f"## File: {rel_path}\n\n"

                # Process each file based on its extension
                file_extension = (
                    file_path.rsplit(".", maxsplit=1)[-1].lower()
                    if "." in file_path
                    else ""
                )
                file_result = None

                try:
                    # Use the same processing logic as process_input
                    if file_extension == "py":
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None, text_content=f.read()
                            )

                    elif file_extension in [
                        "txt",
                        "md",
                        "sh",
                        "yaml",
                        "yml",
                        "toml",
                        "csv",
                    ]:
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None, text_content=f.read()
                            )

                    elif file_extension in ["jsonld", "json"]:
                        with open(file_path, "r", encoding="utf-8") as f:
                            file_result = DocumentConverterResult(
                                title=None,
                                text_content=json.dumps(
                                    json.load(f), ensure_ascii=False, indent=2
                                ),
                            )

                    elif file_extension in ["xlsx", "xls"]:
                        file_result = XlsxConverter(local_path=file_path)

                    elif file_extension == "pdf":
                        file_result = DocumentConverterResult(
                            title=None,
                            text_content=pdfminer.high_level.extract_text(file_path),
                        )

                    elif file_extension in ["docx", "doc"]:
                        file_result = DocxConverter(local_path=file_path)

                    elif file_extension in ["html", "htm"]:
                        file_result = HtmlConverter(local_path=file_path)

                    elif file_extension in ["pptx", "ppt"]:
                        file_result = PptxConverter(local_path=file_path)

                    elif file_extension in IMAGE_EXTENSIONS:
                        # Media files noted but not processed (no external API for captions)
                        md_content += f"[{file_extension.upper()} file - processing not available without external API]\n\n"
                        continue

                    elif file_extension in AUDIO_EXTENSIONS:
                        # Media files noted but not processed (no external API for captions)
                        md_content += f"[{file_extension.upper()} file - processing not available without external API]\n\n"
                        continue

                    elif file_extension in VIDEO_EXTENSIONS:
                        # Media files noted but not processed (no external API for captions)
                        md_content += f"[{file_extension.upper()} file - processing not available without external API]\n\n"
                        continue

                    elif file_extension == "pdb":
                        md_content += "[PDB file - specialized format]\n\n"
                        continue

                    else:
                        # Try MarkItDown as fallback
                        try:
                            md_tool = MarkItDown(enable_plugins=True)
                            file_result = md_tool.convert(file_path)
                        except Exception:
                            md_content += (
                                f"[Unsupported file type: {file_extension}]\n\n"
                            )
                            continue

                    # Add the processed content
                    if file_result and getattr(file_result, "text_content", None):
                        content = file_result.text_content
                        # Limit length for each file
                        max_len = 50_000
                        if len(content) > max_len:
                            content = content[:max_len] + "\n... [Content truncated]"
                        md_content += f"```\n{content}\n```\n\n"

                except Exception as e:
                    md_content += f"[Error processing file: {str(e)}]\n\n"
                    print(f"Warning: Error processing {rel_path} from ZIP: {e}")

    finally:
        # Clean up temporary directory
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            print(f"Warning: Could not remove temporary directory {temp_dir}: {e}")

    return DocumentConverterResult(
        title="ZIP Archive Contents", text_content=md_content.strip()
    )


def extract_pdf_images(
    pdf_path: str,
    output_dir: Optional[str] = None,
    page_range: Optional[tuple] = None
) -> List[Dict[str, Any]]:
    """
    Extract images from PDF file using PyMuPDF.

    Args:
        pdf_path: Path to PDF file
        output_dir: Directory to save extracted images. If None, creates temp directory
        page_range: Tuple of (start_page, end_page) to extract from specific pages.
                    If None, extracts from all pages. Page numbers are 0-indexed.

    Returns:
        List of dictionaries containing image information:
        [
            {
                "page": int,  # Page number (0-indexed)
                "index": int,  # Image index on the page
                "xref": int,  # PDF object reference
                "width": int,  # Image width in pixels
                "height": int,  # Image height in pixels
                "format": str,  # Image format (png, jpeg, etc.)
                "size": int,  # File size in bytes
                "saved_path": str,  # Path where image was saved
            },
            ...
        ]

    Raises:
        ImportError: If PyMuPDF (fitz) is not installed
        FileNotFoundError: If PDF file doesn't exist
    """
    if fitz is None:
        raise ImportError(
            "PyMuPDF (fitz) is required for image extraction. "
            "Install with: pip install pymupdf"
        )

    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF file not found: {pdf_path}")

    # Create output directory if needed
    if output_dir is None:
        output_dir = tempfile.mkdtemp(prefix="pdf_images_")
    else:
        os.makedirs(output_dir, exist_ok=True)

    logger = logging.getLogger(__name__)
    logger.info(f"Extracting images from PDF: {pdf_path}")
    logger.info(f"Saving images to: {output_dir}")

    images_info = []

    try:
        doc = fitz.open(pdf_path)

        # Determine page range
        start_page = 0
        end_page = len(doc)
        if page_range:
            start_page, end_page = page_range
            start_page = max(0, start_page)
            end_page = min(len(doc), end_page)

        # Extract images from each page
        for page_num in range(start_page, end_page):
            page = doc[page_num]
            image_list = page.get_images()

            logger.debug(f"Page {page_num}: found {len(image_list)} images")

            for img_index, img in enumerate(image_list):
                xref = img[0]  # XREF number

                # Extract image
                try:
                    base_image = doc.extract_image(xref)

                    # Get image properties
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]  # png, jpeg, etc.
                    image_width = base_image["width"]
                    image_height = base_image["height"]

                    # Generate filename
                    image_filename = f"page{page_num:03d}_img{img_index:02d}.{image_ext}"
                    image_path = os.path.join(output_dir, image_filename)

                    # Save image
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    # Record image info
                    images_info.append({
                        "page": page_num,
                        "index": img_index,
                        "xref": xref,
                        "width": image_width,
                        "height": image_height,
                        "format": image_ext,
                        "size": len(image_bytes),
                        "saved_path": image_path,
                    })

                    logger.debug(
                        f"Extracted: {image_filename} "
                        f"({image_width}x{image_height}, {len(image_bytes)} bytes)"
                    )

                except Exception as e:
                    logger.warning(
                        f"Failed to extract image {img_index} from page {page_num}: {e}"
                    )
                    continue

        doc.close()

        logger.info(f"Successfully extracted {len(images_info)} images")

    except Exception as e:
        logger.error(f"Error extracting images from PDF: {e}")
        raise

    return images_info


def PdfConverter(
    local_path: str,
    extract_metadata: bool = False,
    extract_sections: bool = False,
    extract_images: bool = False,
    images_output_dir: Optional[str] = None,
    fix_latex: bool = True
) -> DocumentConverterResult:
    """
    Convert a PDF file to text format with enhanced features.

    Args:
        local_path: Path to PDF file to convert.
        extract_metadata: Whether to extract metadata (file size, etc.)
        extract_sections: Whether to extract sections from content
        extract_images: Whether to extract images from PDF (requires PyMuPDF)
        images_output_dir: Directory to save extracted images (default: temp directory)
        fix_latex: Whether to fix LaTeX formula parsing issues

    Returns:
        DocumentConverterResult containing extracted text and optional metadata/sections/images.
    """
    logger = logging.getLogger(__name__)

    if pdfminer is None:
        return DocumentConverterResult(
            title=None,
            text_content="[Error: pdfminer-six not installed]",
            error="pdfminer-six not installed"
        )

    try:
        # Extract text content
        text_content = pdfminer.high_level.extract_text(local_path)

        # Get PDF page count
        pdf_page_count = None
        try:
            from pdfminer.pdfparser import PDFParser
            from pdfminer.pdfdocument import PDFDocument
            from pdfminer.pdfpage import PDFPage

            with open(local_path, 'rb') as fp:
                parser = PDFParser(fp)
                document = PDFDocument(parser)
                pdf_page_count = len(list(PDFPage.create_pages(document)))
        except Exception as e:
            logger.warning(f"Could not extract PDF page count: {e}")

        # Apply LaTeX fixes if requested
        if fix_latex:
            text_content = fix_latex_formulas(text_content)

        # Apply content limit (200,000 characters)
        text_content = apply_content_limit(text_content)

        # Prepare metadata
        metadata = {}
        sections = []
        images = []

        if extract_metadata:
            metadata = {
                "file_path": local_path,
                "file_size": os.path.getsize(local_path) if os.path.exists(local_path) else None,
                "file_extension": os.path.splitext(local_path)[1],
                "conversion_timestamp": time.time(),
                "pdf_page_count": pdf_page_count,  # Add actual PDF page count
            }

        if extract_sections:
            sections = extract_sections_from_markdown(text_content)

        # Extract images if requested
        if extract_images:
            if fitz is None:
                logger.warning(
                    "PyMuPDF not installed. Cannot extract images. "
                    "Install with: pip install pymupdf"
                )
                metadata["image_extraction_error"] = "PyMuPDF not installed"
            else:
                try:
                    images = extract_pdf_images(local_path, output_dir=images_output_dir)
                    logger.info(f"Extracted {len(images)} images from PDF")

                    # Add image count to metadata
                    if extract_metadata:
                        metadata["image_count"] = len(images)
                        if images_output_dir:
                            metadata["images_directory"] = images_output_dir

                except Exception as e:
                    logger.error(f"Failed to extract images: {e}")
                    metadata["image_extraction_error"] = str(e)

        # Try to extract title from first line or metadata
        title = None
        if text_content:
            first_line = text_content.split('\n')[0].strip()
            if first_line and len(first_line) < 200:  # Reasonable title length
                title = first_line

        return DocumentConverterResult(
            title=title,
            text_content=text_content,
            metadata=metadata,
            sections=sections,
            tables=[],  # PDF tables not extracted in basic version
            images=images,  # Extracted images
            processing_time_ms=None  # Can be calculated by caller
        )

    except Exception as e:
        logger.error(f"Error converting PDF: {e}")
        return DocumentConverterResult(
            title=None,
            text_content=f"Error converting PDF: {str(e)}",
            error=str(e)
        )


def TextConverter(local_path: str) -> DocumentConverterResult:
    """
    Read a text file.

    Args:
        local_path: Path to text file to read.

    Returns:
        DocumentConverterResult containing text content.
    """
    with open(local_path, "r", encoding="utf-8") as f:
        text_content = f.read()
    return DocumentConverterResult(title=None, text_content=text_content)


def JsonConverter(local_path: str) -> DocumentConverterResult:
    """
    Read and format a JSON file.

    Args:
        local_path: Path to JSON file to read.

    Returns:
        DocumentConverterResult containing formatted JSON.
    """
    with open(local_path, "r", encoding="utf-8") as f:
        text_content = json.dumps(
            json.load(f), ensure_ascii=False, indent=2
        )
    return DocumentConverterResult(title=None, text_content=text_content)


def YamlConverter(local_path: str) -> DocumentConverterResult:
    """
    Read a YAML file.

    Args:
        local_path: Path to YAML file to read.

    Returns:
        DocumentConverterResult containing YAML content.
    """
    if yaml is None:
        return DocumentConverterResult(
            title=None,
            text_content="[Error: pyyaml not installed]"
        )

    with open(local_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
        text_content = yaml.dump(data, allow_unicode=True, default_flow_style=False)
    return DocumentConverterResult(title=None, text_content=text_content)


def CsvConverter(local_path: str) -> DocumentConverterResult:
    """
    Convert a CSV file to markdown table format.

    Args:
        local_path: Path to CSV file to convert.

    Returns:
        DocumentConverterResult containing markdown table.
    """
    with open(local_path, "r", encoding="utf-8", newline="") as f:
        reader = csv_module.reader(f)
        rows = list(reader)

    if not rows:
        return DocumentConverterResult(title=None, text_content="Empty CSV file")

    md_content = ""
    for i, row in enumerate(rows):
        md_content += "| " + " | ".join(row) + " |\n"
        if i == 0:  # Add a separator after header
            md_content += "|" + "---|" * len(row) + "\n"

    return DocumentConverterResult(title=None, text_content=md_content)


def MarkItDownConverter(local_path: str) -> DocumentConverterResult:
    """
    Convert a file using MarkItDown library (universal converter).

    Args:
        local_path: Path to file to convert.

    Returns:
        DocumentConverterResult containing converted markdown.
    """
    if MarkItDown is None:
        return DocumentConverterResult(
            title=None,
            text_content="[Error: markitdown not installed]"
        )

    md = MarkItDown(enable_plugins=True)
    result = md.convert(local_path)
    return DocumentConverterResult(title=result.title, text_content=result.text_content)


# ============================================================================
# Pagination and Session Management
# ============================================================================

import hashlib
import time
from typing import Optional, Dict, Any, Tuple


class PaginationManager:
    """Manages pagination and session state for large documents."""

    def __init__(self, content: str, page_size: int = 10000):
        """
        Initialize pagination manager.

        Args:
            content: The full content to paginate
            page_size: Number of characters per page (default: 10000)
        """
        self.content = content
        self.page_size = page_size
        self.total_chars = len(content)
        self.total_pages = max(1, (self.total_chars + page_size - 1) // page_size)

    def get_page(self, page: int = 1) -> Tuple[str, bool, Dict[str, Any]]:
        """
        Get a specific page of content.

        Args:
            page: Page number (1-indexed)

        Returns:
            Tuple of (page_content, has_more, pagination_info)
        """
        if page < 1:
            page = 1
        if page > self.total_pages:
            page = self.total_pages

        start = (page - 1) * self.page_size
        end = min(start + self.page_size, self.total_chars)

        page_content = self.content[start:end]
        has_more = end < self.total_chars

        pagination_info = {
            "current_page": page,
            "total_pages": self.total_pages,
            "page_size": self.page_size,
            "char_start": start,
            "char_end": end,
            "has_more": has_more,
            "total_chars": self.total_chars
        }

        return page_content, has_more, pagination_info

    def get_slice(self, offset: int, limit: Optional[int] = None) -> Tuple[str, bool, Dict[str, Any]]:
        """
        Get a slice of content by character offset.

        Args:
            offset: Character offset to start from
            limit: Maximum number of characters to return (None for all remaining)

        Returns:
            Tuple of (slice_content, has_more, pagination_info)
        """
        if offset >= self.total_chars:
            return "", False, {"char_offset": offset, "char_limit": limit, "has_more": False}

        if limit is None:
            end = self.total_chars
            has_more = False
        else:
            end = min(offset + limit, self.total_chars)
            has_more = end < self.total_chars

        slice_content = self.content[offset:end]

        pagination_info = {
            "char_offset": offset,
            "char_limit": limit,
            "char_start": offset,
            "char_end": end,
            "has_more": has_more,
            "total_chars": self.total_chars
        }

        return slice_content, has_more, pagination_info


def generate_session_id(file_path: str, prefix: str = "session") -> str:
    """
    Generate a unique session ID for a file.

    Args:
        file_path: Path to the file
        prefix: Prefix for the session ID

    Returns:
        Unique session ID string
    """
    file_hash = hashlib.md5(file_path.encode()).hexdigest()[:8]
    timestamp = int(time.time())
    return f"{prefix}_{file_hash}_{timestamp}"


def apply_content_limit(content: str, max_chars: int = 200000) -> str:
    """
    Apply hard limit to content length.

    Args:
        content: Content to limit
        max_chars: Maximum number of characters (default: 200,000)

    Returns:
        Limited content with truncation notice if needed
    """
    if len(content) > max_chars:
        return content[:max_chars] + "\n... [Content truncated]"
    return content


def extract_sections_from_markdown(content: str) -> List[Dict[str, Any]]:
    """
    Extract sections from markdown text based on headings.

    Args:
        content: Markdown content

    Returns:
        List of section dictionaries with heading, level, content, etc.
    """
    sections = []
    lines = content.split('\n')
    current_section = None
    section_content = []

    for i, line in enumerate(lines):
        line_stripped = line.strip()
        if line_stripped.startswith('#'):
            # Save previous section if exists
            if current_section is not None:
                sections.append({
                    "heading": current_section['heading'],
                    "level": current_section['level'],
                    "content": '\n'.join(section_content).strip(),
                    "start_line": current_section['start_line'],
                    "end_line": i - 1,
                    "char_start": current_section.get('char_start', 0),
                    "char_end": current_section.get('char_end', 0)
                })

            # Start new section
            heading_text = line.lstrip('#').strip()
            level = len(line) - len(line.lstrip('#'))
            current_section = {
                'heading': heading_text,
                'level': level,
                'start_line': i,
                'char_start': sum(len(l) + 1 for l in lines[:i])  # +1 for newline
            }
            section_content = []
        elif current_section is not None:
            section_content.append(line)

    # Add the last section if exists
    if current_section is not None:
        sections.append({
            "heading": current_section['heading'],
            "level": current_section['level'],
            "content": '\n'.join(section_content).strip(),
            "start_line": current_section['start_line'],
            "end_line": len(lines) - 1,
            "char_start": current_section.get('char_start', 0),
            "char_end": sum(len(line) + 1 for line in lines) - 1
        })

    return sections


def fix_latex_formulas(content: str) -> str:
    """Fix common LaTeX formula parsing issues from PDF extraction.

    This function replaces common LaTeX parsing artifacts with proper Unicode characters:
    - CID placeholders (cid:XXX) with corresponding characters
    - LaTeX commands like \alpha with Greek letters (α)
    - Mathematical symbols like \times with Unicode (×)
    - Simplifies superscripts and subscripts notation

    Args:
        content: Content string with LaTeX formulas that need fixing

    Returns:
        Content with fixed LaTeX formulas converted to Unicode

    Example:
        >>> content = "Formula (cid:16)x(cid:17)"
        >>> result = fix_latex_formulas(content)
        >>> "〈x〉" in result
        True

    Note:
        This is a best-effort conversion. Some complex LaTeX formulas
        may not be fully converted.
    """
    if not content:
        return content

    # Fix (cid:XXX) placeholders - using simple replace instead of regex
    cid_map = {
        '(cid:16)': '〈',
        '(cid:17)': '〉',
        '(cid:40)': '(',
        '(cid:41)': ')',
        '(cid:91)': '[',
        '(cid:93)': ']',
        '(cid:123)': '{',
        '(cid:125)': '}',
        '(cid:60)': '<',
        '(cid:62)': '>',
        '(cid:34)': '"',
        '(cid:39)': "'",
        '(cid:44)': ',',
        '(cid:46)': '.',
        '(cid:58)': ':',
        '(cid:59)': ';',
        '(cid:61)': '=',
        '(cid:43)': '+',
        '(cid:45)': '-',
        '(cid:42)': '*',
        '(cid:47)': '/',
        '(cid:92)': '\\',
        '(cid:124)': '|',
    }
    for pattern, replacement in cid_map.items():
        content = content.replace(pattern, replacement)

    # Fix Greek letters - use replace instead of re.sub
    greek_map = {
        r'\alpha': 'α',
        r'\beta': 'β',
        r'\gamma': 'γ',
        r'\delta': 'δ',
        r'\epsilon': 'ε',
        r'\zeta': 'ζ',
        r'\eta': 'η',
        r'\theta': 'θ',
        r'\iota': 'ι',
        r'\kappa': 'κ',
        r'\lambda': 'λ',
        r'\mu': 'μ',
        r'\nu': 'ν',
        r'\xi': 'ξ',
        r'\pi': 'π',
        r'\rho': 'ρ',
        r'\sigma': 'σ',
        r'\tau': 'τ',
        r'\upsilon': 'υ',
        r'\phi': 'φ',
        r'\chi': 'χ',
        r'\psi': 'ψ',
        r'\omega': 'ω',
        r'\Alpha': 'Α',
        r'\Beta': 'Β',
        r'\Gamma': 'Γ',
        r'\Delta': 'Δ',
        r'\Epsilon': 'Ε',
        r'\Zeta': 'Ζ',
        r'\Eta': 'Η',
        r'\Theta': 'Θ',
        r'\Iota': 'Ι',
        r'\Kappa': 'Κ',
        r'\Lambda': 'Λ',
        r'\Mu': 'Μ',
        r'\Nu': 'Ν',
        r'\Xi': 'Ξ',
        r'\Pi': 'Π',
        r'\Rho': 'Ρ',
        r'\Sigma': 'Σ',
        r'\Tau': 'Τ',
        r'\Upsilon': 'Υ',
        r'\Phi': 'Φ',
        r'\Chi': 'Χ',
        r'\Psi': 'Ψ',
        r'\Omega': 'Ω',
    }
    for latex_cmd, unicode_char in greek_map.items():
        content = content.replace(latex_cmd, unicode_char)

    # Fix mathematical symbols
    math_map = {
        r'\times': '×',
        r'\div': '÷',
        r'\pm': '±',
        r'\mp': '∓',
        r'\leq': '≤',
        r'\geq': '≥',
        r'\neq': '≠',
        r'\approx': '≈',
        r'\equiv': '≡',
        r'\propto': '∝',
        r'\infty': '∞',
        r'\partial': '∂',
        r'\nabla': '∇',
        r'\cdot': '·',
        r'\cdots': '⋯',
        r'\vdots': '⋮',
        r'\ddots': '⋱',
        r'\int': '∫',
        r'\sum': '∑',
        r'\prod': '∏',
        r'\cup': '∪',
        r'\cap': '∩',
        r'\in': '∈',
        r'\notin': '∉',
        r'\subset': '⊂',
        r'\supset': '⊃',
        r'\subseteq': '⊆',
        r'\supseteq': '⊇',
        r'\emptyset': '∅',
        r'\forall': '∀',
        r'\exists': '∃',
        r'\neg': '¬',
        r'\wedge': '∧',
        r'\vee': '∨',
        r'\rightarrow': '→',
        r'\leftarrow': '←',
        r'\Rightarrow': '⇒',
        r'\Leftarrow': '⇐',
        r'\Leftrightarrow': '⇔',
    }
    for latex_cmd, unicode_char in math_map.items():
        content = content.replace(latex_cmd, unicode_char)

    # Fix superscripts and subscripts
    content = re.sub(r'\^\{(\d+)\}', r'^\1', content)  # ^{2} → ^2
    content = re.sub(r'_\{(\d+)\}', r'_\1', content)  # _{2} → _2
    content = re.sub(r'\^\{([a-zA-Z])\}', r'^\1', content)  # ^{x} → ^x
    content = re.sub(r'_\{([a-zA-Z])\}', r'_\1', content)  # _{x} → _x

    return content


    return content
